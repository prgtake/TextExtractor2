# -*- coding: utf-8 -*-
# =====================================================
#  TextExtractor2 v1.1.0
#  Copyright (c) 2026 Datan (データン)
#  Licensed under the MIT License.
# =====================================================
import tkinter as tk
from tkinter import filedialog, messagebox
import tkinter.simpledialog as sd
import os
import json
import re
import sqlite3
import datetime
import time
import csv
import pandas as pd
import docx
from pptx import Presentation
import extract_msg
from email import policy
from email.parser import BytesParser
from google import genai
from google.genai import types
from google.genai.errors import APIError

# =====================================================
#  設定とAPI初期化
# =====================================================
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
CONFIG_FILE = os.path.join(BASE_DIR, "app_config.json")
DEFAULT_MODEL = "gemini-2.5-flash-lite"

# アプリバージョン
APP_VERSION = "1.1.0"

def load_config():
    config = {"GEMINI_API_KEY": os.environ.get("GEMINI_API_KEY"), "MODEL_NAME": DEFAULT_MODEL}
    if os.path.exists(CONFIG_FILE):
        try:
            with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                config.update(json.load(f))
        except:
            pass
    return config

def save_config(api_key, model_name):
    try:
        with open(CONFIG_FILE, "w", encoding="utf-8") as f:
            json.dump({"GEMINI_API_KEY": api_key, "MODEL_NAME": model_name}, f)
    except Exception as e:
        print(f"設定保存失敗: {e}")

class MultiFileGeminiExecutor:
    def __init__(self, master):
        self.master = master
        master.title(f"TextExtractor2 v{APP_VERSION}（サブフォルダ単位抽出）")
        self.master.geometry("600x520") # ボタンサイズに合わせて微調整
        self.master.configure(bg="#f5f5f5")

        self.config = load_config()
        self.api_key = self.config.get("GEMINI_API_KEY")
        
        if not self.api_key or self.api_key == "YOUR_API_KEY_HERE":
            self.api_key = sd.askstring(f"TextExtractor2 v{APP_VERSION}", "Gemini API キーを入力してください：", parent=master)
            if not self.api_key:
                messagebox.showerror(f"TextExtractor2 v{APP_VERSION}", "APIキーが必要です。終了します。")
                master.destroy()
                return
            save_config(self.api_key, self.config.get("MODEL_NAME"))

        try:
            self.client = genai.Client(
                api_key=self.api_key,
                http_options={'retry_options': {'attempts': 5}}
            )
        except Exception as e:
            messagebox.showerror(f"TextExtractor2 v{APP_VERSION}", f"クライアント初期化失敗: {e}")
            self.client = None

        self.target_folder = None
        self.prompt_text = None
        self.prompt_file_path = None
        self.prompt_name = None
        self.columns = []

        self.setup_ui()

    def setup_ui(self):
        # ボタンの幅を統一（19文字分に調整）
        btn_w = 19
        btn_style = {"font": ("Meiryo", 9), "width": btn_w, "pady": 2}
        label_style = {"bg": "#f5f5f5", "font": ("Meiryo", 9)}
        sakura = "#FFB7C5"
        leaf = "#99ff99"

        # 1. フォルダ選択
        f1 = tk.Frame(self.master, bg="#f5f5f5")
        f1.pack(anchor="w", padx=20, pady=(10, 0))
        tk.Button(f1, text="処理対象フォルダを指定", command=self.select_folder, **btn_style).pack(side="left")
        self.folder_label = tk.Label(f1, text="未指定", fg="blue", **label_style)
        self.folder_label.pack(side="left", padx=10)

        # 2. 指示文選択
        f2 = tk.Frame(self.master, bg="#f5f5f5")
        f2.pack(anchor="w", padx=20, pady=5)
        tk.Button(f2, text="指示ファイル(txt)を指定", command=self.select_prompt_file, **btn_style).pack(side="left")
        tk.Button(f2, text="書き方ヘルプ", command=self.show_prompt_help, bg="#ffffcc", **btn_style).pack(side="left", padx=5)
        tk.Button(f2, text="AIと指示文相談", command=self.open_consultation, bg="#e1f5fe", **btn_style).pack(side="left", padx=5)
        self.prompt_label = tk.Label(f2, text="未指定", fg="blue", **label_style)
        self.prompt_label.pack(side="left", padx=5)

        # 3. モデル設定
        f3 = tk.Frame(self.master, bg="#f5f5f5")
        f3.pack(anchor="w", padx=20, pady=5)
        tk.Label(f3, text="使用モデル:", **label_style).pack(side="left")
        self.model_var = tk.StringVar(value=self.config.get("MODEL_NAME", DEFAULT_MODEL))
        tk.Entry(f3, textvariable=self.model_var, width=25).pack(side="left", padx=5)

        # 4. オプション
        opt_frame = tk.LabelFrame(self.master, text="実行オプション", bg="#f5f5f5", padx=10, pady=5)
        opt_frame.pack(fill="x", padx=20, pady=5)

        self.use_web_search = tk.BooleanVar(value=False)
        tk.Checkbutton(opt_frame, text="WEB検索を利用", variable=self.use_web_search, **label_style).grid(row=0, column=0, sticky="w")

        self.use_file_search = tk.BooleanVar(value=False)
        tk.Checkbutton(opt_frame, text="RAG(File Search)を利用", variable=self.use_file_search, **label_style).grid(row=0, column=1, sticky="w", padx=20)

        self.use_sqlite_tool = tk.BooleanVar(value=False)
        tk.Checkbutton(opt_frame, text="SQLiteツールを利用", variable=self.use_sqlite_tool, **label_style).grid(row=1, column=0, sticky="w")

        self.skip_processed = tk.BooleanVar(value=True)
        tk.Checkbutton(opt_frame, text="処理済みファイルをスキップ", variable=self.skip_processed, **label_style).grid(row=1, column=1, sticky="w", padx=20)

        # 5. RAG選択
        rag_f = tk.Frame(self.master, bg="#f5f5f5")
        rag_f.pack(anchor="w", padx=20, pady=2)
        tk.Label(rag_f, text="使用するRAG:", **label_style).pack(side="left")
        self.file_search_store = tk.StringVar(value="")
        self.file_search_dropdown = tk.OptionMenu(rag_f, self.file_search_store, "")
        self.file_search_dropdown.pack(side="left", padx=5)
        self.update_rag_list()

        # 6. 実行ボタン
        btn_f = tk.Frame(self.master, bg="#f5f5f5")
        btn_f.pack(anchor="w", padx=20, pady=10)
        self.run_btn = tk.Button(btn_f, text="指示文を実行", bg=sakura, command=self.run_prompt, **btn_style)
        self.run_btn.pack(side="left")
        tk.Button(btn_f, text="DBをCSV出力", bg=leaf, command=self.export_db_to_csv, **btn_style).pack(side="left", padx=5)

        # 7. ログエリア
        log_f = tk.Frame(self.master, bg="#f5f5f5")
        log_f.pack(fill="both", expand=True, padx=20, pady=5)
        tk.Label(log_f, text="実行ログ:", **label_style).pack(anchor="w")
        self.log_area = tk.Text(log_f, height=6, font=("Consolas", 9), state="disabled", bg="white")
        self.log_area.pack(side="left", fill="both", expand=True)
        sb = tk.Scrollbar(log_f, command=self.log_area.yview)
        sb.pack(side="right", fill="y")
        self.log_area.config(yscrollcommand=sb.set)

        # ログのタグ設定
        self.log_area.tag_config("error", foreground="#d32f2f") # Red
        self.log_area.tag_config("debug", foreground="#757575") # Gray
        self.log_area.tag_config("info", foreground="#1976d2")  # Blue
        self.log_area.tag_config("success", foreground="#388e3c") # Green
        self.log_area.tag_config("warning", foreground="#f57c00") # Orange

    def log_message(self, msg, level="INFO"):
        now = datetime.datetime.now().strftime("%H:%M:%S")
        level = level.upper()
        line = f"[{now}] [{level}] {msg}\n"
        
        tag = level.lower()
        self.log_area.config(state="normal")
        self.log_area.insert("end", line, tag)
        self.log_area.see("end")
        self.log_area.config(state="disabled")
        self.master.update()

    def update_rag_list(self):
        if not self.client: return
        try:
            stores = self.client.file_search_stores.list()
            menu = self.file_search_dropdown["menu"]
            menu.delete(0, "end")
            for store in stores:
                store_id = store.name
                menu.add_command(label=store.display_name or store_id, command=lambda x=store_id: self.file_search_store.set(x))
        except:
            pass

    def show_prompt_help(self):
        help_text = (
            "【指示ファイル（プロンプト）の書き方・重要ルール】\n\n"
            "※文字コードは「UTF-8 (BOMなし)」で保存してください。※\n\n"
            "以下の3つのセクションで構成してください。\n\n"
            "① 概要（任意）\n"
            "   AIに「これから何を解析させるか」を伝えます。\n\n"
            "② 出力項目（必須・1行で記述）\n"
            "   形式： 出力項目: 項目1, 項目2, 項目3\n\n"
            "③ 処理指示（必須・項目ごとに指定）\n"
            "   各項目に対し、以下のタグを使い分けて指示を書きます。\n"
            "   ・【抽出】: 資料からそのまま抜き出す（推測禁止）。\n"
            "   ・【生成】: 資料を元にAIが考える、またはツールを使う。\n\n"
            "【指示ファイルの記述例】\n"
            "----------------------------------------\n"
            "これは取引先からの見積書と提案書のセットです。\n\n"
            "出力項目: 会社名, 合計金額, 提案の要約, 印象評価\n\n"
            "### 処理指示\n"
            "【会社名】：【抽出】見積書に記載されている発行元の会社名。\n"
            "【合計金額】：【抽出】見積書の税込み合計金額。数値のみ。\n"
            "【提案の要約】：【生成】提案書の内容を100文字以内で要約。\n"
            "【印象評価】：【生成】提案内容が革新的か保守的か判定してください。\n"
            "----------------------------------------\n\n"
            "【高精度に計算させるコツ】\n"
            "合計点や判定を求める場合は、必ずその「根拠となる数値（各科目の点数など）」も出力項目に含め、先に書き出させるようにしてください。これだけで計算精度が劇的に向上します。"
        )
        messagebox.showinfo(f"TextExtractor2 v{APP_VERSION}", help_text)

    def open_consultation(self):
        if not self.client:
            messagebox.showerror(f"TextExtractor2 v{APP_VERSION}", "APIクライアントが初期化されていません。")
            return
        PromptConsultationWindow(self.master, self.client)

    def select_folder(self):
        path = filedialog.askdirectory()
        if path:
            self.target_folder = path
            self.folder_label.config(text=path)

    def select_prompt_file(self):
        path = filedialog.askopenfilename(filetypes=[("テキスト", "*.txt")])
        if not path: return
        self.prompt_file_path = path
        self.prompt_name = os.path.splitext(os.path.basename(path))[0]
        self.prompt_label.config(text=os.path.basename(path))
        with open(path, "r", encoding="utf-8") as f:
            self.prompt_text = f.read()
        m = re.search(r"出力項目[:：]\s*(.+)", self.prompt_text)
        if not m:
            messagebox.showerror(f"TextExtractor2 v{APP_VERSION}", "プロンプトに『出力項目:』が見つかりません。")
            self.columns = []
            return
        self.columns = [c.strip() for c in re.split(r'[,、]', m.group(1))]

    def extract_text_from_any_file(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()
        content = ""
        try:
            if ext in [".txt", ".csv", ".log"]:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
            elif ext in [".xlsx", ".xls"]:
                dict_df = pd.read_excel(file_path, sheet_name=None)
                for sheet, df in dict_df.items():
                    content += f"\n[Sheet: {sheet}]\n{df.to_csv(index=False)}"
            elif ext == ".docx":
                doc = docx.Document(file_path)
                content = "\n".join([p.text for p in doc.paragraphs])
                for table in doc.tables:
                    for row in table.rows:
                        content += "\n" + "|".join([cell.text for cell in row.cells])
            elif ext == ".pptx":
                prs = Presentation(file_path)
                content = "\n".join([shape.text for s in prs.slides for shape in s.shapes if hasattr(shape, "text")])
            elif ext in [".doc", ".ppt", ".jtd"]:
                import olefile
                if olefile.isOleFile(file_path):
                    with olefile.OleFileIO(file_path) as ole:
                        for s_list in ole.listdir():
                            s_path = "/".join(s_list)
                            if any(k in s_path for k in ["WordDocument", "PowerPoint Document", "JS-Main", "Body"]):
                                with ole.openstream(s_list) as stream:
                                    data = stream.read()
                                    try:
                                        content += data.decode('cp932', errors='ignore')
                                    except:
                                        content += "".join([chr(b) if 32 <= b <= 126 or b > 128 else " " for b in data])
            elif ext == ".rtf":
                from striprtf.striprtf import rtf_to_text
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = rtf_to_text(f.read())
            elif ext == ".odt":
                from odf import text, teletype, opendocument
                textdoc = opendocument.load(file_path)
                content = "\n".join([teletype.extractText(p) for p in textdoc.getElementsByType(text.P)])
            elif ext == ".eml":
                with open(file_path, 'rb') as f:
                    msg = BytesParser(policy=policy.default).parse(f)
                content = f"Subject: {msg['subject']}\nFrom: {msg['from']}\nDate: {msg['date']}\n\n"
                content += msg.get_body(preferencelist=('plain')).get_content()
            elif ext == ".msg":
                msg = extract_msg.Message(file_path)
                content = f"Subject: {msg.subject}\nFrom: {msg.sender}\nDate: {msg.date}\n\n{msg.body}"
                msg.close()

            if content.strip():
                self.log_message(f"テキスト抽出成功 ({ext}): {os.path.basename(file_path)} ({len(content)}文字)", "DEBUG")
            else:
                self.log_message(f"テキスト抽出結果が空です ({ext}): {os.path.basename(file_path)}", "WARNING")

        except Exception as e:
            self.log_message(f"抽出失敗 ({ext}): {e}", "ERROR")
        return content

    def execute_sqlite_tool(self, db_path: str, sql_query: str) -> str:
        self.log_message(f"SQLiteツール実行: {sql_query[:100]}...", "DEBUG")
        query_upper = sql_query.strip().upper()
        # 参照系（SELECT, PRAGMA）以外の実行を厳格に禁止
        is_read_only = query_upper.startswith("SELECT") or query_upper.startswith("PRAGMA")
        
        if not is_read_only:
            self.log_message(f"更新系SQLがブロックされました: {sql_query[:50]}", "WARNING")
            return "Error: This tool is strictly read-only. Only SELECT or PRAGMA queries are allowed."

        try:
            conn = sqlite3.connect(db_path)
            cur = conn.cursor()
            if hasattr(self, 'current_sqlite_paths'):
                for path in self.current_sqlite_paths:
                    if os.path.abspath(path) != os.path.abspath(db_path) and os.path.exists(path):
                        alias = os.path.splitext(os.path.basename(path))[0]
                        cur.execute(f"ATTACH DATABASE '{path}' AS {alias}")
            
            cur.execute(sql_query)
            rows = cur.fetchall()
            col_names = [d[0] for d in cur.description] if cur.description else []
            res = [dict(zip(col_names, r)) for r in rows]
            conn.close()
            self.log_message(f"SQL取得成功: {len(rows)}件", "DEBUG")
            return json.dumps(res, ensure_ascii=False)
            
        except Exception as e:
            self.log_message(f"SQL実行エラー: {e}", "ERROR")
            return f"Error: {str(e)}"

    def get_system_instruction(self):
        instruction = (
            f"\n\n### システムルール (厳守) ###\n"
            f"1. 出力は必ず以下の項目をキーとするJSON配列形式 [{{...}}, {{...}}] のみを返してください。\n"
            f"   出力項目: {', '.join(self.columns)}\n"
            f"2. 各項目について、【抽出】（ファイル内容のみ使用、検索禁止）と【生成】（ツールや知見を使用）を適切に使い分けてください。\n"
            f"3. 内容が見つからない項目は、推測せず必ず空文字 (\"\") としてください。\n"
            f"4. 解説、挨拶、Markdownの装飾(```json等)は一切禁止。パース可能な純粋なJSONデータのみを返してください。\n"
            f"5. 文字化けや不自然な文字列は、文脈から日本語として意味が通じるよう再解釈してください。\n"
            f"6. ファイルの「フルパス」や「保存日時」の情報と「ファイル自体の記述内容」が矛盾する場合（例：フォルダ名と書類内の氏名が異なる等）は、必ず【ファイル自体の記述内容】を真実として優先してください。\n"
        )
        if self.use_sqlite_tool.get():
            instruction += "\n7. SQLite連携が有効です。ATTACHは自動実行済みです。「エイリアス名.テーブル名」で指定してください。\n"
        return instruction

    def run_prompt(self):
        if not self.target_folder or not self.prompt_text or not self.columns:
            messagebox.showerror(f"TextExtractor2 v{APP_VERSION}", "設定が不足しています。")
            return
        model_name = self.model_var.get()
        save_config(self.api_key, model_name)

        sqlite_paths = []
        if self.use_sqlite_tool.get():
            tag = "###使用するSQLiteのパス"
            if tag not in self.prompt_text:
                messagebox.showerror(f"TextExtractor2 v{APP_VERSION}", f"SQLite使用にチェックがありますがプロンプト内に「{tag}」がありません。")
                return
            sqlite_paths = re.findall(r'"([^"]+\.db(?:sqlite\d?)?)"', self.prompt_text)
            self.current_sqlite_paths = sqlite_paths
            for p in sqlite_paths:
                if not os.path.exists(p):
                    messagebox.showerror(f"TextExtractor2 v{APP_VERSION}", f"SQLiteファイルが見つかりません: {p}")
                    return

        self.run_btn.config(state="disabled")
        self.log_area.config(state="normal")
        self.log_area.delete("1.0", "end")
        self.log_area.config(state="disabled")
        self.log_message(f"=== 処理開始 (Model: {model_name}) ===")
        self.log_message(f"対象フォルダ: {self.target_folder}")

        conn = self.setup_db()
        cur = conn.cursor()

        try:
            subfolders = [os.path.join(self.target_folder, d) for d in os.listdir(self.target_folder) if os.path.isdir(os.path.join(self.target_folder, d))]
            if not subfolders: 
                subfolders = [self.target_folder]
                self.log_message("サブフォルダが見つかりません。対象フォルダ自体を1つの単位として処理します。", "DEBUG")

            self.log_message(f"合計処理単位数: {len(subfolders)}")

            system_instr = self.get_system_instruction()

            for i, sub in enumerate(subfolders):
                sub_name = os.path.basename(sub)
                files_in_sub = []
                for root, _, files in os.walk(sub):
                    for f in files:
                        ext = os.path.splitext(f)[1].lower()
                        if ext in [".db", ".exe", ".zip", ".py", ".txt"]: continue
                        files_in_sub.append(os.path.join(root, f))

                if not files_in_sub:
                    self.log_message(f"[{i+1}/{len(subfolders)}] フォルダ空につきスキップ: {sub_name}", "DEBUG")
                    continue

                self.log_message(f"[{i+1}/{len(subfolders)}] 解析中: {sub_name} ({len(files_in_sub)}ファイル)")
                if self.skip_processed.get():
                    cur.execute("SELECT COUNT(*) FROM extracted_data WHERE folder_path = ?", (sub,))
                    if cur.fetchone()[0] > 0:
                        self.log_message(f"  -> スキップ (処理済み記録あり)")
                        continue

                contents = [f"{self.prompt_text}\n{system_instr}\n\n※サブフォルダ内の複数ファイルを統合して解析してください。\n"]
                parts = []
                for f_path in files_in_sub:
                    ext = os.path.splitext(f_path)[1].lower()
                    full_path = os.path.abspath(f_path)
                    mtime = os.path.getmtime(f_path)
                    save_time = datetime.datetime.fromtimestamp(mtime).strftime("%Y-%m-%d %H:%M:%S")

                    if ext in [".pdf", ".png", ".jpg", ".jpeg", ".webp"]:
                        self.log_message(f"  -> バイナリ添付: {os.path.basename(f_path)}", "DEBUG")
                        contents[0] += f"\n--- バイナリファイル: {full_path} (保存日時: {save_time}) ---\n"
                        with open(f_path, "rb") as f: fb = f.read()
                        mime = "application/pdf" if ext == ".pdf" else "image/jpeg"
                        parts.append(types.Part.from_bytes(data=fb, mime_type=mime))
                    else:
                        txt = self.extract_text_from_any_file(f_path)
                        if txt.strip():
                            contents[0] += f"\n--- ファイル: {full_path} (保存日時: {save_time}) ---\n{txt}\n"

                items = self.call_gemini_api(contents + parts, model_name, sqlite_paths)
                if items:
                    for item in items:
                        values = [sub, sub_name, "SUBFOLDER_UNIT", *[str(item.get(col, "")) for col in self.columns], datetime.datetime.now().isoformat()]
                        cur.execute(f"INSERT INTO extracted_data (folder_path, folder_name, file_path, {','.join(['\"'+c+'\"' for c in self.columns])}, created_at) "
                                    f"VALUES ({','.join(['?']*len(values))})", values)
                    conn.commit()
                    self.log_message(f"  -> 成功 ({len(items)}件のデータを保存)", "SUCCESS")
                else:
                    self.log_message(f"  -> 失敗 (このフォルダの解析結果が得られませんでした)", "ERROR")

        except Exception as main_e:
            self.log_message(f"実行中に致命的なエラーが発生しました: {main_e}", "ERROR")

        finally:
            conn.close()
            self.run_btn.config(state="normal")
            self.log_message("=== すべての処理が完了しました ===")
            messagebox.showinfo(f"TextExtractor2 v{APP_VERSION}", "すべての処理が終了しました。")

    def _execute_api_call(self, contents, config, model_name, purpose="API"):
        # contents: List of parts (str or types.Part) or a single string
        if isinstance(contents, str):
            contents = [contents]
        
        # ツール呼び出しのループ（最大10回）
        messages = [types.Content(role="user", parts=[
            types.Part.from_text(text=p) if isinstance(p, str) else p for p in contents
        ])]

        for turn in range(10):
            max_retries = 5
            response = None
            for attempt in range(max_retries):
                try:
                    self.log_message(f"[{purpose}] APIリクエスト中 (試行 {attempt+1}/{max_retries}, ターン {turn+1})...", "DEBUG")
                    response = self.client.models.generate_content(
                        model=model_name,
                        contents=messages,
                        config=config
                    )
                    break
                except Exception as e:
                    err_msg = str(e)
                    is_retryable = any(x in err_msg for x in ["429", "503", "500", "504", "UNAVAILABLE", "Resource has been exhausted"])
                    if is_retryable and attempt < max_retries - 1:
                        wait = 30 * (2 ** attempt)
                        self.log_message(f"一時エラー ({err_msg})。{wait}秒待機して再試行...", "WARNING")
                        time.sleep(wait)
                    else:
                        self.log_message(f"APIエラー: {e}", "ERROR")
                        return None

            if not response or not response.candidates:
                return None

            # トークン使用量の記録
            if hasattr(response, 'usage_metadata'):
                meta = response.usage_metadata
                self.log_message(f"トークン使用量: 入力={meta.prompt_token_count}, 出力={meta.candidates_token_count}", "DEBUG")

            cand = response.candidates[0]
            if not cand.content or not cand.content.parts:
                return response

            messages.append(cand.content)

            # ツール呼び出し（Function Calling）の抽出
            tool_calls = [p.function_call for p in cand.content.parts if p.function_call]
            if not tool_calls:
                return response

            # ツールの実行とレスポンスの生成
            tool_responses = []
            for call in tool_calls:
                if call.name == "execute_sqlite_tool":
                    sql = call.args.get("sql_query")
                    # 第一パス目に使用されるDBパス（通常はプロンプト内で指定されたもの）
                    target_db = self.current_sqlite_paths[0] if hasattr(self, 'current_sqlite_paths') and self.current_sqlite_paths else ""
                    res_val = self.execute_sqlite_tool(target_db, sql)
                    tool_responses.append(types.Part.from_function_response(
                        name=call.name, response={"result": res_val}
                    ))
                else:
                    tool_responses.append(types.Part.from_function_response(
                        name=call.name, response={"error": f"Unknown tool: {call.name}"}
                    ))

            if tool_responses:
                messages.append(types.Content(role="user", parts=tool_responses))
            else:
                return response
        return None

    def call_gemini_api(self, contents, model_name, sqlite_paths):
        # contents[0] はシステム指示含むプロンプト、残りはバイナリパーツ
        original_prompt = contents[0]
        binary_parts = contents[1:]
        
        has_web = self.use_web_search.get()
        has_fss = self.use_file_search.get() and self.file_search_store.get()
        has_sqlite = self.use_sqlite_tool.get() and sqlite_paths

        # ツールが1つ以下なら、効率のため単一リクエストで実行
        active_tools_count = sum([bool(has_web), bool(has_fss), bool(bool(has_sqlite))])

        if active_tools_count <= 1:
            config = {"system_instruction": self.get_system_instruction(sqlite_paths), "response_mime_type": "application/json"}
            if has_web:
                config["tools"] = [{"google_search": {}}]
                del config["response_mime_type"]
            elif has_fss:
                config["tools"] = [{"file_search": {"file_search_store_names": [self.file_search_store.get()]}}]
                del config["response_mime_type"]
            elif has_sqlite:
                config["tools"] = [self.execute_sqlite_tool]
                del config["response_mime_type"]
            
            res = self._execute_api_call(contents, config, model_name, purpose="AI直接解析")
            return self._parse_response(res)

        # 複数ツールがある場合は、逐次リサーチステージ（制限回避）
        self.log_message("複数ツールが選択されたため、逐次リサーチを開始します。", "INFO")
        combined_research = ""

        # 1. WEB検索
        if has_web:
            cfg = {"tools": [{"google_search": {}}]}
            prompt = f"以下の指示に基づいて、WEB検索を行い必要な情報を調査してください。回答は調査結果のまとめのみを返してください。\n\n{original_prompt}"
            res = self._execute_api_call(prompt, cfg, model_name, purpose="WEB検索")
            if res and res.text:
                combined_research += f"### WEB検索結果:\n{res.text}\n\n"

        # 2. RAG(File Search)
        if has_fss:
            cfg = {"tools": [{"file_search": {"file_search_store_names": [self.file_search_store.get()]}}]}
            prompt = f"以下の調査成果と指示に基づき、RAG(File Search)を用いて追加情報を調査してください。\n\n【これまでの調査成果】\n{combined_research}\n\n【元の指示】\n{original_prompt}"
            res = self._execute_api_call(prompt, cfg, model_name, purpose="RAG検索")
            if res and res.text:
                combined_research += f"### RAG(File Search)結果:\n{res.text}\n\n"

        # 3. SQLiteツール
        if has_sqlite:
            # SQLiteの場合はシステムプロンプトのスキーマ情報が必要なため get_system_instruction を考慮
            cfg = {"system_instruction": self.get_system_instruction(), "tools": [self.execute_sqlite_tool]}
            db_prompt = f"以下の調査成果と指示に基づき、SQLiteツールを用いて必要なデータを取得してください。\n\n【これまでの調査成果】\n{combined_research}\n\n【元の指示】\n{original_prompt}"
            res = self._execute_api_call(db_prompt, cfg, model_name, purpose="DB検索")
            if res and res.text:
                combined_research += f"### SQLiteツール結果:\n{res.text}\n\n"

        # 4. 最終統合：全ての成果を統合してJSON出力
        self.log_message("最終的なデータ統合とJSON生成を実行中...", "INFO")
        final_prompt = (
            f"これまでの全ての調査結果と添付されたファイルを統合し、最終的な抽出結果をJSON形式で出力してください。\n\n"
            f"【調査結果まとめ】\n{combined_research}\n\n"
            f"【元の指示とシステムルール】\n{original_prompt}"
        )
        # ツールなしの最終段階で JSON Mode を使用
        final_config = {"response_mime_type": "application/json"}
        res = self._execute_api_call([final_prompt] + binary_parts, final_config, model_name, purpose="最終統合")
        return self._parse_response(res)

    def _parse_response(self, response):
        if not response or not hasattr(response, 'text'):
            return None
        raw_text = response.text.strip()
        match = re.search(r"(\[.*\]|\{.*\})", raw_text, re.DOTALL)
        if not match:
            self.log_message("API応答からJSONが見つかりませんでした。", "ERROR")
            self.log_message(raw_text[:500], "DEBUG")
            return None
        try:
            parsed = json.loads(match.group(1))
            return parsed if isinstance(parsed, list) else [parsed]
        except Exception as e:
            self.log_message(f"JSONパースエラー: {e}", "ERROR")
            return None

    def setup_db(self):
        db_path = os.path.join(self.target_folder, f"{self.prompt_name}_DB.db")
        conn = sqlite3.connect(db_path)
        cur = conn.cursor()
        cur.execute(f"CREATE TABLE IF NOT EXISTS extracted_data (id INTEGER PRIMARY KEY AUTOINCREMENT, folder_path TEXT, folder_name TEXT, file_path TEXT, {', '.join([f'\"{c}\" TEXT' for c in self.columns])}, created_at TEXT)")
        conn.commit()
        return conn

    def export_db_to_csv(self):
        if not self.target_folder or not self.prompt_name:
            messagebox.showerror(f"TextExtractor2 v{APP_VERSION}", "設定が不足しています。")
            return
        db_path = os.path.join(self.target_folder, f"{self.prompt_name}_DB.db")
        if not os.path.exists(db_path):
            messagebox.showerror(f"TextExtractor2 v{APP_VERSION}", "DBがありません。")
            return
        csv_path = filedialog.asksaveasfilename(defaultextension=".csv", filetypes=[("CSV", "*.csv")])
        if not csv_path: return
        try:
            conn = sqlite3.connect(db_path)
            df = pd.read_sql("SELECT * FROM extracted_data", conn)
            df = df.applymap(lambda x: x.replace("\r\n", " ").replace("\n", " ").replace("\r", " ") if isinstance(x, str) else x)
            df.to_csv(csv_path, index=False, encoding="utf_8_sig", quoting=csv.QUOTE_ALL)
            conn.close()
            messagebox.showinfo(f"TextExtractor2 v{APP_VERSION}", "CSV出力完了（セル内改行はスペースに変換されました）")
        except Exception as e:
            messagebox.showerror(f"TextExtractor2 v{APP_VERSION}", f"CSV出力に失敗しました: {e}")

class PromptConsultationWindow:
    def __init__(self, parent, client):
        self.window = tk.Toplevel(parent)
        self.window.title(f"TextExtractor2 v{APP_VERSION} - AI指示文相談チャット")
        self.window.geometry("500x600")
        self.client = client

        # 自らのソースコードを読み込む
        try:
            import sys
            import os
            if hasattr(sys, '_MEIPASS'):
                source_path = os.path.join(sys._MEIPASS, "textextractor2.py")
            else:
                source_path = os.path.abspath(__file__)
                
            with open(source_path, "r", encoding="utf-8") as f:
                source_code = f.read()
        except Exception as e:
            source_code = f"ソースコードの読み込みに失敗しました: {e}"

        # システム指示（プロンプト職人としての設定）
        self.system_instruction = (
            "あなたは『TextExtractor2』というツールの【指示文作成アドバイザー】です。\n"
            "このツールのソースコードを以下に提供します。これを読み込み、仕様を完璧に理解してください：\n"
            f"--- SOURCE CODE START ---\n{source_code}\n--- SOURCE CODE END ---\n\n"
            "ユーザーの要望を聞き、最適な【指示ファイル（.txt）】の内容を提案してください。\n"
            "【重要ルール】\n"
            "1. 回答は常に【簡潔・最小限】にしてください。余計な解説は省き、必要なことだけを返してください。\n"
            "2. 指示文案（概要、出力項目、処理指示）は、ユーザーがそのままコピペできるよう、必ず「```text ... ```」のようなコードブロック形式で出力してください。\n"
            "3. 【生成】項目については、どのように生成・推論するか具体的な基準を確認してください。\n"
            "4. WEB検索、RAG、SQLiteツールの使用有無を必ず確認してください。\n"
            "5. SQLiteを使用する場合は、使うDB의パスを確認し、実行時に「エイリアス名.テーブル名」という形式を使用するルールを徹底してください。"
        )
        self.chat = []
        self.txt_area = tk.Text(self.window, state="disabled", font=("Meiryo", 9), wrap="word", bg="#fafafa")
        self.txt_area.pack(fill="both", expand=True, padx=10, pady=10)
        self.entry = tk.Entry(self.window, font=("Meiryo", 10))
        self.entry.pack(fill="x", padx=10, pady=(0, 10))
        self.entry.bind("<Return>", lambda e: self.send_message())
        btn = tk.Button(self.window, text="送信 (Enter)", command=self.send_message, bg="#e1f5fe")
        btn.pack(pady=(0, 10))
        self.display_message("AI", (
            "こんにちは！指示文（プロンプト）の作成をサポートします。\n"
            "指示文の典型的な書き方は、次のようなものです。\n\n"
            "【例：フォルダ内の見積書と提案書を解析する場合】\n\n"
            "⇒ 指示文の例：\n"
            "--------------------------------------------\n"
            "###使用するSQLiteのパス\n"
            "\"C:\\data\\MasterData.db\"\n\n"
            "出力項目: 会社名, 合計金額, 提案の要約, 顧客ランク\n\n"
            "### 処理指示\n"
            "・【会社名】：【抽出】見積書に記載されている発行元の会社名。\n"
            "・【合計金額】：【抽出】見積書の税込み合計金額。数値のみ。\n"
            "・【提案の要約】：【生成】提案書の内容を100文字以内で要約。\n"
            "・【顧客ランク】：【生成】MasterData.顧客マスタを参照し、【会社名】の取引ランクを取得。\n"
            "--------------------------------------------\n\n"
            "※項目名を【項目名】と書くことで、AIが他の抽出・生成結果を参照して計算や照合に利用できます。\n"
            "※外部DB参照時は必ず「エイリアス名.テーブル名」の形式で指定してください。\n\n"
            "今回は、どのような抽出やデータ生成を行いたいですか？"
        ))

    def display_message(self, sender, text):
        self.txt_area.config(state="normal")
        self.txt_area.insert("end", f"【{sender}】\n{text}\n\n")
        self.txt_area.see("end")
        self.txt_area.config(state="disabled")

    def send_message(self):
        user_text = self.entry.get()
        if not user_text.strip(): return
        self.entry.delete(0, "end")
        self.display_message("あなた", user_text)
        self.chat.append({"role": "user", "parts": [user_text]})
        try:
            response = self.client.models.generate_content(
                model="gemini-2.5-flash-lite",
                contents=[self.system_instruction] + [types.Content(role=c["role"], parts=[types.Part.from_text(text=c["parts"][0])]) for c in self.chat]
            )
            ai_text = response.text
            self.display_message("AI", ai_text)
            self.chat.append({"role": "model", "parts": [ai_text]})
        except Exception as e:
            self.display_message("Error", str(e))

if __name__ == "__main__":
    root = tk.Tk()
    app = MultiFileGeminiExecutor(root)
    root.mainloop()
